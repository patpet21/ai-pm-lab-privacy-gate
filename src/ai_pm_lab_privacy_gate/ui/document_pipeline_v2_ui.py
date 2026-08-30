from __future__ import annotations

import html
import os
from pathlib import Path
from types import MethodType

from pptx import Presentation
from PySide6.QtCore import QTimer, Qt
from PySide6.QtGui import QImageReader, QPixmap
from PySide6.QtWidgets import (
    QFileDialog,
    QFrame,
    QLabel,
    QMessageBox,
    QScrollArea,
    QTextBrowser,
)


def _pptx_slide_texts(path: Path) -> list[tuple[str, str]]:
    presentation = Presentation(str(path))
    slides: list[tuple[str, str]] = []
    for index, slide in enumerate(presentation.slides, start=1):
        lines: list[str] = []
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False):
                for paragraph in shape.text_frame.paragraphs:
                    if paragraph.text.strip():
                        lines.append(paragraph.text)
            if getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    values = [cell.text.strip() for cell in row.cells]
                    if any(values):
                        lines.append(" | ".join(values))
        slides.append((f"Slide {index}", "\n".join(lines)))
    return slides


def _load_pptx_internal(view, path: Path, protected: bool) -> None:
    view.clear()
    for title, text in _pptx_slide_texts(path):
        browser = QTextBrowser()
        browser.setObjectName("OfficeDocumentView")
        escaped = html.escape(text).replace("\n", "<br>") or "&nbsp;"
        if protected:
            escaped = view._highlight(escaped, True)
        browser.setHtml(
            "<html><body style='font-family:Segoe UI;color:#102A43;background:#fff;margin:20px'>"
            f"<h3>{html.escape(title)}</h3><p>{escaped}</p></body></html>"
        )
        view.tabs.addTab(browser, title)


def _image_preview_widget() -> tuple[QScrollArea, QLabel]:
    """Reuse the existing document panel and add a neutral raster viewport."""
    label = QLabel()
    label.setAlignment(Qt.AlignmentFlag.AlignCenter)
    label.setScaledContents(False)
    scroll = QScrollArea()
    scroll.setFrameShape(QFrame.Shape.NoFrame)
    scroll.setWidgetResizable(True)
    scroll.setWidget(label)
    return scroll, label


def _load_image_preview(label: QLabel, path: Path) -> None:
    reader = QImageReader(str(path))
    reader.setAutoTransform(True)
    image = reader.read()
    if image.isNull():
        detail = reader.errorString().strip() or "unknown image decode error"
        raise ValueError(f"Image preview unavailable: {detail}")
    pixmap = QPixmap.fromImage(image)
    # Preview-only downscaling protects the UI from very large phone photos.
    # Protection always runs against the original full-resolution pixels.
    if pixmap.width() > 1800 or pixmap.height() > 1800:
        pixmap = pixmap.scaled(
            1800,
            1800,
            Qt.AspectRatioMode.KeepAspectRatio,
            Qt.TransformationMode.SmoothTransformation,
        )
    label.setPixmap(pixmap)


def _apply_restore_pipeline_v2(main_window) -> None:
    restore = getattr(main_window, "restore_page", None)
    if restore is None or getattr(restore, "_document_pipeline_v2_applied", False):
        return

    restore.FILE_FILTER = (
        "Supported results (*.txt *.pdf *.docx *.xlsx *.pptx);;"
        "Text files (*.txt);;PDF files (*.pdf);;Word files (*.docx);;"
        "Excel files (*.xlsx);;PowerPoint files (*.pptx)"
    )
    restore.drop_zone.SUPPORTED = {".txt", ".pdf", ".docx", ".xlsx", ".pptx"}
    restore.drop_zone.formats.setText("TXT, PDF, Word, Excel, PowerPoint — processed locally")

    original_file_loaded = restore._file_loaded

    def file_loaded(self, payload: object) -> None:
        original_file_loaded(payload)
        if self._source_path is not None and self._source_path.suffix.lower() == ".pptx":
            self.preview_tabs.setTabVisible(1, True)
            self.preview_tabs.setCurrentIndex(1)
            self._refresh_document_preview()

    restore._file_loaded = MethodType(file_loaded, restore)

    original_restore_ready = restore._restore_ready

    def restore_ready(self, payload: object) -> None:
        original_restore_ready(payload)
        if self._restored_path is not None and self._restored_path.suffix.lower() == ".pptx":
            self.download_button.setVisible(True)
            self.download_button.setText("Download restored PowerPoint file")
            self.preview_tabs.setTabVisible(1, True)
            self.preview_tabs.setCurrentIndex(1)
            self._refresh_document_preview()

    restore._restore_ready = MethodType(restore_ready, restore)

    original_refresh_document_preview = restore._refresh_document_preview

    def refresh_document_preview(self) -> None:
        if self._source_path is None or self._source_path.suffix.lower() != ".pptx":
            original_refresh_document_preview()
            return
        try:
            self.input_pdf_document.close()
            self.output_pdf_document.close()
            self.high_fidelity_button.setVisible(self._libreoffice_available)
            self.install_libreoffice_button.setVisible(not self._libreoffice_available)
            if self._libreoffice_available and self.high_fidelity_button.isChecked():
                self._set_pdf_controls(True)
                input_pdf = self.office_renderer.render(
                    self._source_path, self._preview_root / "restore-input-office"
                )
                self.input_view_stack.setCurrentIndex(0)
                self.output_view_stack.setCurrentIndex(0)
                self.input_pdf_document.load(str(input_pdf))
                if self._restored_path:
                    output_pdf = self.office_renderer.render(
                        self._restored_path, self._preview_root / "restore-output-office"
                    )
                    self.output_pdf_document.load(str(output_pdf))
                self.preview_note.setText(
                    "High-fidelity local PowerPoint preview: protected result on the left, restored deck on the right."
                )
                self._set_pdf_page(0)
            else:
                self.high_fidelity_button.setChecked(False)
                self.libreoffice_note.setText(
                    "Built-in slide preview active. Install LibreOffice for page-accurate rendering."
                )
                self._set_pdf_controls(False)
                _load_pptx_internal(self.input_office_view, self._source_path, True)
                self.input_view_stack.setCurrentIndex(1)
                self.output_office_view.clear()
                if self._restored_path:
                    _load_pptx_internal(self.output_office_view, self._restored_path, False)
                self.output_view_stack.setCurrentIndex(1)
                self.input_office_view.synchronize_with(self.output_office_view)
                self.preview_note.setText(
                    "Built-in PowerPoint preview: protected result on the left, locally restored deck on the right."
                )
        except Exception as exc:
            self.preview_note.setText(f"PowerPoint preview unavailable: {exc}")

    restore._refresh_document_preview = MethodType(refresh_document_preview, restore)
    try:
        restore.high_fidelity_button.toggled.disconnect()
    except (RuntimeError, TypeError):
        pass
    restore.high_fidelity_button.toggled.connect(lambda _checked: restore._refresh_document_preview())
    restore._document_pipeline_v2_applied = True


def apply_document_pipeline_v2_ui(main_window) -> None:
    """Expose every unified pipeline format without duplicating Protect UI."""
    page = getattr(main_window, "protection_page", None)
    if page is None or getattr(page, "_document_pipeline_v2_applied", False):
        _apply_restore_pipeline_v2(main_window)
        return

    page.pdf_path.setPlaceholderText("PDF, Word, Excel, PowerPoint, TXT, PNG or JPG")
    page.high_fidelity_button.setToolTip(
        "Use LibreOffice locally for a page-accurate Word, Excel or PowerPoint preview."
    )
    page.save_download_button.setToolTip(
        "Save locally and export the protected source-format copy plus a protected TXT companion."
    )

    # The base document panels are shared by PDF/Office previews. Add image
    # viewports to those same stacks instead of creating a second preview system.
    page.original_image_scroll, page.original_image_label = _image_preview_widget()
    page.protected_image_scroll, page.protected_image_label = _image_preview_widget()
    page.original_view_stack.addWidget(page.original_image_scroll)
    page.protected_view_stack.addWidget(page.protected_image_scroll)

    def refresh_initial_image_preview(path_text: str) -> None:
        source = Path(path_text.strip())
        if source.suffix.lower() not in {".png", ".jpg", ".jpeg"} or not source.is_file():
            return
        try:
            _load_image_preview(page.original_image_label, source)
            page.protected_image_label.clear()
            page.protected_image_label.setText("Protected preview available after Scan & Protect")
            page.original_view_stack.setCurrentWidget(page.original_image_scroll)
            page.protected_view_stack.setCurrentWidget(page.protected_image_scroll)
            page.preview_tabs.setTabVisible(1, True)
            page.preview_tabs.setCurrentIndex(1)
            page.comparison_note.setText(
                "Image loaded locally. The original is shown on the left; scan to create "
                "the pixel-redacted protected preview on the right."
            )
        except Exception as exc:
            page.preview_tabs.setTabToolTip(1, f"Preview unavailable: {exc}")
            page.comparison_note.setText(str(exc))

    def queue_initial_image_preview(path_text: str) -> None:
        # Other shared source-state handlers run on the same signal. Queue the
        # image preview once so their reset work completes before it is shown.
        QTimer.singleShot(
            0,
            lambda value=path_text: refresh_initial_image_preview(value),
        )

    page.pdf_path.textChanged.connect(queue_initial_image_preview)

    def browse_document(self) -> None:
        path, _ = QFileDialog.getOpenFileName(
            self,
            "Choose a document",
            "",
            "Supported documents (*.pdf *.docx *.xlsx *.pptx *.txt *.png *.jpg *.jpeg);;"
            "PDF files (*.pdf);;Word files (*.docx);;Excel files (*.xlsx);;"
            "PowerPoint files (*.pptx);;Text files (*.txt);;"
            "Image files (*.png *.jpg *.jpeg)",
        )
        if path:
            self.pdf_path.setText(path)
            self.input_tabs.setCurrentIndex(1)

    try:
        page.browse_button.clicked.disconnect()
    except (RuntimeError, TypeError):
        pass
    page._browse_document = MethodType(browse_document, page)
    page.browse_button.clicked.connect(page._browse_document)

    original_analysis_ready = page._analysis_ready

    def analysis_ready(self, payload: object) -> None:
        original_analysis_ready(payload)
        if self.current_document is not None and self.current_document.source_kind == "image":
            self.focus_preview_button.setChecked(True)

    page._analysis_ready = MethodType(analysis_ready, page)

    original_refresh = page._refresh_preview

    def refresh_preview(self, *args) -> None:
        original_refresh(*args)
        if self.current_document is None:
            return
        if self.current_document.source_kind in {"pptx", "image"}:
            self.preview_tabs.setTabVisible(1, True)
            self._pdf_preview_timer.start()

    page._refresh_preview = MethodType(refresh_preview, page)

    original_comparison = page._update_document_comparison

    def update_document_comparison(self) -> None:
        if (
            self.current_document is not None
            and self.current_result is not None
            and self.current_document.source_path is not None
            and self.current_document.source_kind == "image"
        ):
            try:
                self.protected_pdf_document.close()
                self.original_pdf_document.close()
                self.office_preview_options_widget.setVisible(False)
                self._set_pdf_controls_enabled(False)
                suffix = self.current_document.source_path.suffix.lower()
                if suffix not in {".png", ".jpg", ".jpeg"}:
                    suffix = ".png"
                protected_path = self._preview_directory / (
                    f"protected-image-{os.getpid()}-{self.current_document.source_path.stem}{suffix}"
                )
                self.service.save_protected_document(
                    self.current_result,
                    protected_path,
                    source_document=self.current_document,
                )
                _load_image_preview(
                    self.original_image_label, Path(self.current_document.source_path)
                )
                _load_image_preview(self.protected_image_label, protected_path)
                self.original_view_stack.setCurrentWidget(self.original_image_scroll)
                self.protected_view_stack.setCurrentWidget(self.protected_image_scroll)
                self.comparison_note.setText(
                    "Local image preview: original on the left and pixel-redacted protected copy "
                    "on the right. OCR and redaction stay on this device; source EXIF metadata is "
                    "not copied to the protected image."
                )
            except Exception as exc:
                self.preview_tabs.setTabToolTip(1, f"Preview unavailable: {exc}")
                self.comparison_note.setText(str(exc))
                self.preview_tabs.setCurrentIndex(0)
            else:
                self.preview_tabs.setTabToolTip(
                    1,
                    "Compare the local image source with its protected pixel-redacted copy.",
                )
                self.preview_tabs.setCurrentIndex(1)
            return

        if (
            self.current_document is None
            or self.current_result is None
            or self.current_document.source_path is None
            or self.current_document.source_kind != "pptx"
        ):
            original_comparison()
            return
        try:
            self.protected_pdf_document.close()
            self.original_pdf_document.close()
            protected_path = self._preview_directory / f"protected-office-{Path(self.current_document.source_path).stem}.pptx"
            self.service.save_protected_document(
                self.current_result,
                protected_path,
                source_document=self.current_document,
            )
            self.office_preview_options_widget.setVisible(True)
            self.high_fidelity_button.setVisible(self._libreoffice_available)
            self.install_libreoffice_button.setVisible(not self._libreoffice_available)
            if self._libreoffice_available and self.high_fidelity_button.isChecked():
                original_path = self._office_preview_renderer.render(
                    self.current_document.source_path, self._office_original_directory
                )
                protected_pdf = self._office_preview_renderer.render(
                    protected_path, self._office_protected_directory
                )
                self.original_view_stack.setCurrentIndex(0)
                self.protected_view_stack.setCurrentIndex(0)
                self._set_pdf_controls_enabled(True)
                self.original_pdf_document.load(str(original_path))
                self.protected_pdf_document.load(str(protected_pdf))
                self.comparison_note.setText(
                    "High-fidelity local PowerPoint rendering: original on the left and protected deck on the right."
                )
                self._set_pdf_page(0)
            else:
                self.high_fidelity_button.setChecked(False)
                self.libreoffice_note.setText(
                    "Built-in slide preview active. Install LibreOffice for page-accurate rendering."
                )
                _load_pptx_internal(self.original_office_view, self.current_document.source_path, False)
                _load_pptx_internal(self.protected_office_view, protected_path, True)
                self.original_office_view.synchronize_with(self.protected_office_view)
                self.original_view_stack.setCurrentIndex(1)
                self.protected_view_stack.setCurrentIndex(1)
                self._set_pdf_controls_enabled(False)
                self.comparison_note.setText(
                    "Built-in PowerPoint preview: original on the left and editable protected deck on the right."
                )
        except Exception as exc:
            self.preview_tabs.setTabToolTip(1, f"Preview unavailable: {exc}")
            self.comparison_note.setText(str(exc))
            self.preview_tabs.setCurrentIndex(0)
        else:
            self.preview_tabs.setTabToolTip(
                1, "Compare the local PowerPoint source with its protected editable copy."
            )
            self.preview_tabs.setCurrentIndex(1)

    try:
        page._pdf_preview_timer.timeout.disconnect()
    except (RuntimeError, TypeError):
        pass
    page._update_document_comparison = MethodType(update_document_comparison, page)
    page._pdf_preview_timer.timeout.connect(page._update_document_comparison)

    def save_and_download(self) -> None:
        if not self._confirm_residual_risk("downloading"):
            return
        document = self._save_to_library()
        if not document or not self.current_result or not self.current_document:
            return

        kind = self.current_document.source_kind
        if (
            kind == "image"
            and self.current_document.source_path is not None
            and self.current_document.source_path.suffix.lower() in {".png", ".jpg", ".jpeg"}
        ):
            label = "Image"
            suffix = self.current_document.source_path.suffix.lower()
        else:
            labels = {
                "pdf": ("PDF", ".pdf"),
                "docx": ("Word", ".docx"),
                "xlsx": ("Excel", ".xlsx"),
                "pptx": ("PowerPoint", ".pptx"),
                "txt": ("Text", ".txt"),
                "text": ("Text", ".txt"),
            }
            label, suffix = labels.get(kind, ("Document", ".txt"))

        suggested = f"{document.title}_protected{suffix}"
        path, _ = QFileDialog.getSaveFileName(
            self,
            f"Save protected {label}",
            suggested,
            f"{label} files (*{suffix})",
        )
        if not path:
            return
        destination = Path(path)
        if destination.suffix.lower() != suffix:
            destination = destination.with_suffix(suffix)

        main_output, companion = self.service.save_protected_bundle(
            self.current_result,
            destination,
            source_document=self.current_document,
        )
        if main_output == companion:
            QMessageBox.information(
                self,
                "Protected file exported",
                f"Protected TXT saved to:\n{main_output}",
            )
            return

        QMessageBox.information(
            self,
            "Protected files exported",
            f"Protected {suffix.upper().lstrip('.')} and TXT companion saved locally:\n"
            f"{main_output}\n{companion}",
        )

    try:
        page.save_download_button.clicked.disconnect()
    except (RuntimeError, TypeError):
        pass
    page._save_and_download = MethodType(save_and_download, page)
    page.save_download_button.clicked.connect(page._save_and_download)

    page._document_pipeline_v2_applied = True
    _apply_restore_pipeline_v2(main_window)
