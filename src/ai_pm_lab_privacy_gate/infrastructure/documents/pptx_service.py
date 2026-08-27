from __future__ import annotations

from dataclasses import dataclass
from pathlib import Path
from typing import Any, Iterator

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from ai_pm_lab_privacy_gate.domain.models import AnalysisDocument, PageContent, ProtectionResult


@dataclass(frozen=True, slots=True)
class _PptxPart:
    paragraph: Any
    location: str


class PowerPointDocumentService:
    """Extract and protect editable PowerPoint text while preserving the deck."""

    SUPPORTED_SUFFIXES = {".pptx"}

    def extract(self, path: str | Path) -> AnalysisDocument:
        source = self._validated_source(path)
        presentation = Presentation(str(source))
        entries = [
            (part.paragraph.text, part.location)
            for part in self._iter_parts(presentation)
            if part.paragraph.text.strip()
        ]
        pages = tuple(
            PageContent(page_number=index, text=text, location=location)
            for index, (text, location) in enumerate(entries, start=1)
        )
        return AnalysisDocument(source_kind="pptx", source_path=source, pages=pages)

    def write_protected(
        self,
        source_document: AnalysisDocument,
        result: ProtectionResult,
        path: str | Path,
    ) -> Path:
        if source_document.source_path is None or source_document.source_kind != "pptx":
            raise ValueError("A PowerPoint source document is required.")
        source = self._validated_source(source_document.source_path)
        destination = Path(path)
        if destination.suffix.lower() != ".pptx":
            destination = destination.with_suffix(".pptx")
        destination.parent.mkdir(parents=True, exist_ok=True)

        presentation = Presentation(str(source))
        parts = [part for part in self._iter_parts(presentation) if part.paragraph.text.strip()]
        current_texts = [part.paragraph.text for part in parts]
        expected = [page.text for page in source_document.pages]
        if current_texts != expected:
            raise ValueError(
                "The PowerPoint changed after it was scanned. Scan it again before exporting."
            )

        replacements = self._replacements_by_page(result)
        for page_number, part in enumerate(parts, start=1):
            for start, end, replacement in reversed(replacements.get(page_number, ())):
                self._replace_paragraph_range(part.paragraph, start, end, replacement)

        properties = presentation.core_properties
        properties.author = "AI PM LAB Privacy Gate"
        properties.last_modified_by = "AI PM LAB Privacy Gate"
        properties.comments = "Protected locally by AI PM LAB Privacy Gate"
        presentation.save(str(destination))
        return destination

    @classmethod
    def _validated_source(cls, path: str | Path) -> Path:
        source = Path(path)
        if not source.exists():
            raise FileNotFoundError(source)
        if source.suffix.lower() not in cls.SUPPORTED_SUFFIXES:
            raise ValueError("Supported PowerPoint format is .pptx.")
        return source

    @classmethod
    def _iter_parts(cls, presentation) -> Iterator[_PptxPart]:
        for slide_number, slide in enumerate(presentation.slides, start=1):
            yield from cls._iter_shapes(slide.shapes, slide_number)
            try:
                notes_frame = slide.notes_slide.notes_text_frame
            except (AttributeError, ValueError):
                notes_frame = None
            if notes_frame is not None:
                for index, paragraph in enumerate(notes_frame.paragraphs, start=1):
                    yield _PptxPart(paragraph, f"Slide {slide_number} · Notes {index}")

    @classmethod
    def _iter_shapes(cls, shapes, slide_number: int, prefix: str = "") -> Iterator[_PptxPart]:
        for shape_index, shape in enumerate(shapes, start=1):
            shape_name = str(getattr(shape, "name", "") or f"Shape {shape_index}")
            location = f"Slide {slide_number} · {prefix}{shape_name}"
            if getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.GROUP:
                yield from cls._iter_shapes(
                    shape.shapes,
                    slide_number,
                    prefix=f"{prefix}{shape_name} / ",
                )
                continue
            if getattr(shape, "has_text_frame", False):
                for paragraph_index, paragraph in enumerate(shape.text_frame.paragraphs, start=1):
                    yield _PptxPart(paragraph, f"{location} · Paragraph {paragraph_index}")
            if getattr(shape, "has_table", False):
                for row_index, row in enumerate(shape.table.rows, start=1):
                    for column_index, cell in enumerate(row.cells, start=1):
                        for paragraph_index, paragraph in enumerate(cell.text_frame.paragraphs, start=1):
                            yield _PptxPart(
                                paragraph,
                                f"{location} · Table R{row_index}C{column_index} · Paragraph {paragraph_index}",
                            )

    @staticmethod
    def _replacements_by_page(
        result: ProtectionResult,
    ) -> dict[int, list[tuple[int, int, str]]]:
        replacement_by_finding = {
            span.finding_id: span.replacement_text for span in result.protected_spans
        }
        by_page: dict[int, list[tuple[int, int, str]]] = {}
        for finding in result.applied_findings:
            replacement = replacement_by_finding.get(finding.finding_id)
            if replacement is not None:
                by_page.setdefault(finding.page_number, []).append(
                    (finding.start, finding.end, replacement)
                )
        return by_page

    @staticmethod
    def _replace_paragraph_range(
        paragraph: Any, start: int, end: int, replacement: str
    ) -> None:
        runs = list(paragraph.runs)
        if not runs or "".join(run.text for run in runs) != paragraph.text:
            text = paragraph.text
            paragraph.text = text[:start] + replacement + text[end:]
            return

        cursor = 0
        start_run = start_offset = end_run = end_offset = None
        for index, run in enumerate(runs):
            next_cursor = cursor + len(run.text)
            if start_run is None and (
                start < next_cursor or (start == next_cursor and index == len(runs) - 1)
            ):
                start_run, start_offset = index, start - cursor
            if end_run is None and end <= next_cursor:
                end_run, end_offset = index, end - cursor
                break
            cursor = next_cursor

        if any(value is None for value in (start_run, start_offset, end_run, end_offset)):
            raise ValueError("Unable to map a protected value back to the PowerPoint deck.")
        assert start_run is not None and start_offset is not None
        assert end_run is not None and end_offset is not None
        if start_run == end_run:
            run = runs[start_run]
            run.text = run.text[:start_offset] + replacement + run.text[end_offset:]
            return
        runs[start_run].text = runs[start_run].text[:start_offset] + replacement
        for index in range(start_run + 1, end_run):
            runs[index].text = ""
        runs[end_run].text = runs[end_run].text[end_offset:]
