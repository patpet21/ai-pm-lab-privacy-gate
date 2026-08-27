from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches

from ai_pm_lab_privacy_gate.application.privacy_service import PrivacyGateService
from ai_pm_lab_privacy_gate.domain.models import Finding
from ai_pm_lab_privacy_gate.infrastructure.connectors.gmail_import import _collect_attachments
from ai_pm_lab_privacy_gate.infrastructure.connectors.google_drive_import import _GOOGLE_EXPORTS
from ai_pm_lab_privacy_gate.infrastructure.documents.restore_service import DocumentRestoreService


def _finding(document, value: str, entity_type: str = "PERSON") -> Finding:
    for page in document.pages:
        start = page.text.find(value)
        if start >= 0:
            return Finding(
                finding_id=f"test-{page.page_number}-{start}",
                entity_type=entity_type,
                text=value,
                start=start,
                end=start + len(value),
                score=1.0,
                page_number=page.page_number,
                context=page.text,
            )
    raise AssertionError(f"{value!r} was not extracted")


def test_txt_uses_same_document_pipeline_and_exports_utf8(tmp_path: Path) -> None:
    source = tmp_path / "message.txt"
    source.write_text("Contact Jane Smith at jane@example.com", encoding="utf-8")
    service = PrivacyGateService()

    document = service.document_from_file(source)
    assert document.source_kind == "txt"
    finding = _finding(document, "Jane Smith")
    result = service.protect(document, (finding,))

    output = service.save_protected_document(result, tmp_path / "message_protected.txt", document)
    assert output.read_text(encoding="utf-8") == result.combined_text
    assert "Jane Smith" not in result.combined_text
    assert "[[PG_PERSON_001]]" in result.combined_text


def test_pptx_protect_companion_and_restore_round_trip(tmp_path: Path) -> None:
    source = tmp_path / "deck.pptx"
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(7), Inches(1))
    box.text_frame.text = "Owner Jane Smith · jane@example.com"
    presentation.save(source)

    service = PrivacyGateService()
    document = service.document_from_file(source)
    assert document.source_kind == "pptx"
    result = service.protect(document, (_finding(document, "Jane Smith"),))

    protected = service.save_protected_document(result, tmp_path / "deck_protected.pptx", document)
    companion = service.save_protected_text(result, tmp_path / "deck_protected.txt")

    protected_text = DocumentRestoreService.extract_text(protected)
    assert "Jane Smith" not in protected_text
    assert "[[PG_PERSON_001]]" in protected_text
    assert companion.read_text(encoding="utf-8") == result.combined_text

    restored = DocumentRestoreService().restore(
        protected,
        result.mappings,
        tmp_path / "deck_restored.pptx",
    )
    assert restored.restored_occurrences == 1
    assert "Jane Smith" in DocumentRestoreService.extract_text(restored.output_path)


def test_google_slides_export_to_pptx() -> None:
    mime, suffix = _GOOGLE_EXPORTS["application/vnd.google-apps.presentation"]
    assert mime == "application/vnd.openxmlformats-officedocument.presentationml.presentation"
    assert suffix == ".pptx"


def test_gmail_supported_attachments_include_pptx_and_txt() -> None:
    payload = {
        "mimeType": "multipart/mixed",
        "parts": [
            {
                "filename": "deck.pptx",
                "mimeType": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
                "partId": "1",
                "body": {"attachmentId": "A1"},
            },
            {
                "filename": "notes.txt",
                "mimeType": "text/plain",
                "partId": "2",
                "body": {"attachmentId": "A2"},
            },
            {
                "filename": "photo.jpg",
                "mimeType": "image/jpeg",
                "partId": "3",
                "body": {"attachmentId": "A3"},
            },
        ],
    }
    attachments = _collect_attachments(payload)
    assert [item.filename for item in attachments] == ["deck.pptx", "notes.txt"]
